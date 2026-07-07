"""Text extraction from uploaded documents — Markdown output via pymupdf4llm."""
import logging
import os
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, List, Dict

logger = logging.getLogger("harven")

# Extensions that are explicitly known to be unsupported (no parser exists),
# mapped to an actionable detail message. Checked before any parser dispatch
# so we never attempt to parse a binary/legacy format with the wrong tool
# (e.g. .doc through the python-docx/.docx parser).
_UNSUPPORTED_EXTENSIONS: Dict[str, str] = {
    ".doc": "Formato .doc não suportado — reenvie como .docx ou .pdf",
}


@dataclass
class ExtractionResult:
    """Structured outcome of a text extraction attempt.

    status:
        - "ok": extraction succeeded and produced non-empty text.
        - "empty": extraction ran without error but yielded no text
          (e.g. a scanned PDF with no text layer).
        - "unsupported": no parser exists for this file type.
        - "failed": a parser ran and raised an exception.
    """

    status: str
    text: Optional[str] = None
    detail: Optional[str] = None


def extract(file_path: str, mime_type: str = "") -> ExtractionResult:
    """Extract text as Markdown from a document file, with a structured result.

    Dispatches by extension/mime type to a dedicated parser branch. Every
    parser call is wrapped so a corrupted/malformed file never propagates
    an exception to the caller (mapped to status="failed" instead).
    """
    ext = Path(file_path).suffix.lower()

    if ext in _UNSUPPORTED_EXTENSIONS:
        detail = _UNSUPPORTED_EXTENSIONS[ext]
        logger.warning(f"Unsupported file type for extraction: {ext} — {detail}")
        return ExtractionResult(status="unsupported", detail=detail)

    if ext == ".pdf" or "pdf" in mime_type:
        parser = _extract_pdf_markdown
    elif ext == ".docx" or "wordprocessingml.document" in mime_type:
        parser = _extract_docx_markdown
    elif ext == ".pptx" or "presentationml.presentation" in mime_type:
        parser = _extract_pptx_markdown
    elif ext in (".txt", ".md", ".html", ".htm", ".csv"):
        parser = _extract_plain
    else:
        detail = f"Extensão '{ext or mime_type or 'desconhecida'}' não suportada para extração de texto"
        logger.warning(f"Unsupported file type for extraction: {ext}")
        return ExtractionResult(status="unsupported", detail=detail)

    try:
        text = parser(file_path)
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
        return ExtractionResult(status="failed", detail=str(e))

    if not text or not text.strip():
        return ExtractionResult(status="empty")

    return ExtractionResult(status="ok", text=text)


def extract_text(file_path: str, mime_type: str = "") -> Optional[str]:
    """Legacy wrapper: extract text as Markdown, returning Optional[str].

    Preserves the pre-existing contract for callers that predate the
    structured ``extract()`` result (returns the text on success, ``None``
    for empty/unsupported/failed).
    """
    result = extract(file_path, mime_type)
    return result.text if result.status == "ok" else None


def extract_text_from_bytes(data: bytes, filename: str, mime_type: str = "") -> Optional[str]:
    """Extract text from in-memory bytes."""
    ext = Path(filename).suffix.lower() or ".bin"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        return extract_text(tmp_path, mime_type)
    finally:
        os.unlink(tmp_path)


def extract_chapters_from_bytes(data: bytes, filename: str, mime_type: str = "") -> List[Dict[str, str]]:
    """Extract text and split into chapters based on headings."""
    ext = Path(filename).suffix.lower() or ".bin"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        md = extract_text(tmp_path, mime_type)
        if not md:
            return []
        return split_markdown_into_chapters(md)
    finally:
        os.unlink(tmp_path)


def _is_heading_line(line: str) -> Optional[str]:
    """Detect various heading patterns. Returns the cleaned title or None."""
    stripped = line.strip()
    if not stripped:
        return None
    # Markdown headings: #, ##, or ###
    m = re.match(r"^#{1,3}\s+(.+)", stripped)
    if m:
        return m.group(1).strip()
    # Numbered sections: "1. Topic", "2. Topic", "1 - Topic"
    m = re.match(r"^\d+[\.\)\-]\s+(.{3,})$", stripped)
    if m:
        return m.group(1).strip()
    # Bold lines used as headers: **Topic Name**
    m = re.match(r"^\*\*(.{3,})\*\*\s*$", stripped)
    if m:
        return m.group(1).strip()
    # Slide-style titles: ALL CAPS, short (3-80 chars), standalone
    if stripped.isupper() and 3 <= len(stripped) <= 80 and not stripped.startswith("|"):
        # Exclude table rows and common artifacts
        if not re.match(r"^[-=|+\s]+$", stripped):
            return stripped.title()
    return None


def split_markdown_into_chapters(md: str) -> List[Dict[str, str]]:
    """Split markdown by headings, numbered sections, bold titles, or ALL-CAPS slide titles."""
    lines = md.split("\n")
    chapters: List[Dict[str, str]] = []
    current_title = ""
    current_lines: list[str] = []

    for line in lines:
        heading = _is_heading_line(line)
        if heading:
            if current_lines:
                body = "\n".join(current_lines).strip()
                if body:
                    chapters.append({"title": current_title or "Introducao", "body": body})
            current_title = heading
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        body = "\n".join(current_lines).strip()
        if body:
            chapters.append({"title": current_title or "Conteudo", "body": body})

    if not chapters and md.strip():
        chapters.append({"title": "Conteudo completo", "body": md.strip()})

    return chapters


def _clean_markdown(md: str) -> str:
    """Clean up pymupdf4llm output for better readability."""
    # ---- Strip HTML artifacts ----
    # Remove literal <br>, <br/>, <br /> tags (very common in PDF extraction)
    md = re.sub(r"<br\s*/?\s*>", "\n", md, flags=re.IGNORECASE)
    # Remove other stray HTML tags
    md = re.sub(r"</?(?:p|div|span|font|b|i|u|em|strong|sup|sub|ol|ul|li|td|tr|th|table|img|a)[^>]*>", "", md, flags=re.IGNORECASE)
    # Decode common HTML entities
    md = md.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    md = md.replace("&nbsp;", " ").replace("&quot;", '"').replace("&#39;", "'")

    # ---- Remove PDF artifacts ----
    # Image placeholders
    md = re.sub(r"^\*\*==>.*?<==\*\*\s*$", "", md, flags=re.MULTILINE)
    md = re.sub(r"^==>.*?<==\s*$", "", md, flags=re.MULTILINE)
    # Standalone numbers on their own line (slide/page numbers)
    md = re.sub(r"^\d{1,3}\s*$", "", md, flags=re.MULTILINE)
    # Common PDF headers/footers
    md = re.sub(r"^(Sumário|SUMÁRIO|sumário|Índice|ÍNDICE)\s*$", "", md, flags=re.MULTILINE)
    md = re.sub(r"^Página\s+\d+\s*(de\s+\d+)?\s*$", "", md, flags=re.MULTILINE | re.IGNORECASE)
    # Lines that are just dashes, equals, or underscores (horizontal rules from tables)
    md = re.sub(r"^[-=_]{4,}\s*$", "---", md, flags=re.MULTILINE)

    # ---- Fix broken tables ----
    # Remove malformed markdown tables (rows with pipes but no content)
    md = re.sub(r"^\|[\s|—–-]*\|\s*$", "", md, flags=re.MULTILINE)
    # Fix double pipes
    md = re.sub(r"\|\|+", "| |", md)
    # Remove lines that are just "| |---|---|---| " style separators with no header before
    md = re.sub(r"^\|\s*[-—–:|\s]+\|\s*$", "", md, flags=re.MULTILINE)

    # ---- Reconstruct text flow ----
    # Join broken sentences: line ending without sentence-end punctuation followed by lowercase
    md = re.sub(r"([a-záàâãéèêíïóôõúüç,;])\s*\n([a-záàâãéèêíïóôõúüç])", r"\1 \2", md)
    # Join lines where previous line ends with a hyphen (word-wrap break)
    md = re.sub(r"(\w)-\s*\n(\w)", r"\1\2", md)

    # ---- Normalize formatting ----
    # Bullet points
    md = re.sub(r"^[•●◦▪▸►]\s*", "- ", md, flags=re.MULTILINE)
    md = re.sub(r"^[–—]\s+", "- ", md, flags=re.MULTILINE)
    # Remove excessive blank lines (3+ → 2)
    md = re.sub(r"\n{3,}", "\n\n", md)
    # Remove excessive whitespace within lines
    md = re.sub(r"[ \t]{3,}", " ", md)
    # Remove trailing whitespace per line
    md = re.sub(r" +$", "", md, flags=re.MULTILINE)

    return md.strip()


def _extract_pdf_markdown(path: str) -> Optional[str]:
    import pymupdf4llm

    md = pymupdf4llm.to_markdown(path)
    if not md or not md.strip():
        return None
    return _clean_markdown(md)


def _extract_docx_markdown(path: str) -> Optional[str]:
    from docx import Document

    doc = Document(path)
    lines = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower()
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        else:
            lines.append(text)
        lines.append("")
    result = "\n".join(lines)
    return result.strip() if result.strip() else None


def _extract_plain(path: str) -> Optional[str]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip() or None


def _extract_pptx_markdown(path: str) -> Optional[str]:
    """Extract text from a .pptx presentation, slide by slide, in order.

    Walks every shape on every slide (including grouped shapes and tables)
    and concatenates any text found, using a heading per slide so the
    output stays consistent with the Markdown produced by the other
    extractors.
    """
    from pptx import Presentation

    prs = Presentation(path)
    lines: List[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_lines: List[str] = []
        for shape in slide.shapes:
            slide_lines.extend(_pptx_shape_text_lines(shape))

        slide_text = [line for line in slide_lines if line.strip()]
        if not slide_text:
            continue

        lines.append(f"## Slide {slide_index}")
        lines.extend(slide_text)
        lines.append("")

    result = "\n".join(lines)
    return result.strip() if result.strip() else None


def _pptx_shape_text_lines(shape) -> List[str]:
    """Recursively collect text lines from a python-pptx shape.

    Handles plain text frames, tables, and grouped shapes (which nest
    child shapes that must be walked the same way).
    """
    lines: List[str] = []

    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            lines.extend(_pptx_shape_text_lines(sub_shape))
        return lines

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if not text:
                text = paragraph.text.strip()
            if text:
                lines.append(text)

    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                lines.append(row_text)

    return lines
