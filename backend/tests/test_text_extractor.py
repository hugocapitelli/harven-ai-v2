"""Tests for services.text_extractor — structured extract() + legacy shim.

Covers FILE-1's 5 required outcomes:
  1. .pptx with text -> status="ok"
  2. .doc -> status="unsupported" (no exception)
  3. Scanned PDF / no text layer -> status="empty"
  4. Corrupted file -> status="failed" (no crash)
  5. extract_text() legacy caller still gets Optional[str]
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from services.text_extractor import ExtractionResult, extract, extract_text


# ---------------------------------------------------------------------------
# Fixtures — real .pptx files built with python-pptx (no mocking of the
# parser itself, so the extraction logic is exercised end-to-end).
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_with_text(tmp_path: Path) -> str:
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Introducao"
    body = slide1.placeholders[1]
    body.text_frame.text = "Primeira linha"
    p = body.text_frame.add_paragraph()
    p.text = "Segunda linha"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(1))
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(0, 1).text = "B"
    table_shape.table.cell(1, 0).text = "C"
    table_shape.table.cell(1, 1).text = "D"

    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def pptx_empty(tmp_path: Path) -> str:
    """A .pptx with a slide but no textual content anywhere."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout, no shapes added
    path = tmp_path / "empty.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def doc_file(tmp_path: Path) -> str:
    """A file with a .doc extension (legacy OLE2 binary — content is irrelevant,
    the extension alone must short-circuit to unsupported)."""
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1garbage-ole2-bytes")
    return str(path)


@pytest.fixture
def corrupted_pptx(tmp_path: Path) -> str:
    """A file with a .pptx extension whose content is not a valid zip/OOXML
    package, so the python-pptx parser raises."""
    path = tmp_path / "corrupt.pptx"
    path.write_text("this is not a real pptx package")
    return str(path)


# ---------------------------------------------------------------------------
# extract() — structured result
# ---------------------------------------------------------------------------


class TestExtractPptx:
    def test_pptx_with_text_returns_ok_with_concatenated_text(self, pptx_with_text: str):
        result = extract(pptx_with_text)

        assert isinstance(result, ExtractionResult)
        assert result.status == "ok"
        assert result.text is not None
        assert result.detail is None

    def test_pptx_text_preserves_slide_order_and_content(self, pptx_with_text: str):
        result = extract(pptx_with_text)

        assert result.text.index("Introducao") < result.text.index("Primeira linha")
        assert result.text.index("Primeira linha") < result.text.index("Segunda linha")
        # Table cells from slide 2 must also be present.
        assert "A" in result.text and "D" in result.text

    def test_pptx_dispatch_by_mime_type_when_extension_missing(self, tmp_path: Path, pptx_with_text: str):
        # Copy the valid pptx bytes to a path without the .pptx suffix, and
        # rely on mime_type for dispatch (mirrors extract_text_from_bytes()
        # which always writes a suffix, but extract() itself must also honor
        # mime_type per the .pdf/.docx branches' existing convention).
        no_ext_path = tmp_path / "deck_no_ext"
        no_ext_path.write_bytes(Path(pptx_with_text).read_bytes())

        result = extract(
            str(no_ext_path),
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        assert result.status == "ok"
        assert "Introducao" in result.text

    def test_pptx_without_text_returns_empty(self, pptx_empty: str):
        result = extract(pptx_empty)

        assert result.status == "empty"
        assert result.text is None


class TestExtractDoc:
    def test_doc_returns_unsupported_with_actionable_detail(self, doc_file: str):
        result = extract(doc_file)

        assert result.status == "unsupported"
        assert result.text is None
        assert result.detail
        assert ".docx" in result.detail or ".pdf" in result.detail

    def test_doc_never_raises(self, doc_file: str):
        # Must not attempt to parse .doc as .docx and must not propagate
        # any exception — this is the core of BUG-SWEEP #9.
        try:
            extract(doc_file)
        except Exception as e:  # pragma: no cover - failure path documented
            pytest.fail(f"extract() raised for a .doc file: {e}")


class TestExtractUnknownExtension:
    def test_unknown_extension_returns_unsupported(self, tmp_path: Path):
        path = tmp_path / "file.xyz"
        path.write_text("hello")

        result = extract(str(path))

        assert result.status == "unsupported"
        assert result.detail


class TestExtractEmptyPdf:
    def test_pdf_with_no_text_layer_returns_empty(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
        # Simulate a scanned PDF (parser runs successfully but yields no text)
        # by mocking pymupdf4llm.to_markdown directly — avoids depending on a
        # real scanned-PDF fixture while still exercising the empty-vs-failed
        # distinction that is the whole point of AC4.
        import pymupdf4llm

        monkeypatch.setattr(pymupdf4llm, "to_markdown", lambda path: "   \n\n  ")

        path = tmp_path / "scanned.pdf"
        path.write_bytes(b"%PDF-1.4 fake")

        result = extract(str(path))

        assert result.status == "empty"
        assert result.text is None


class TestExtractFailure:
    def test_corrupted_pptx_returns_failed_without_raising(self, corrupted_pptx: str):
        result = extract(corrupted_pptx)

        assert result.status == "failed"
        assert result.detail

    def test_parser_exception_is_captured_for_any_format(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch):
        import services.text_extractor as text_extractor

        def _boom(path: str):
            raise RuntimeError("simulated parser crash")

        monkeypatch.setattr(text_extractor, "_extract_docx_markdown", _boom)

        path = tmp_path / "broken.docx"
        path.write_text("not really docx")

        result = extract(str(path))

        assert result.status == "failed"
        assert "simulated parser crash" in result.detail


class TestExtractOkOnlyWhenNonEmpty:
    def test_status_ok_never_returned_with_empty_text(self, tmp_path: Path):
        path = tmp_path / "blank.txt"
        path.write_text("   \n\n   ")

        result = extract(str(path))

        assert result.status == "empty"
        assert result.text is None


# ---------------------------------------------------------------------------
# extract_text() — legacy Optional[str] contract must remain unbroken
# ---------------------------------------------------------------------------


class TestExtractTextLegacyShim:
    def test_returns_text_on_success(self, pptx_with_text: str):
        text = extract_text(pptx_with_text)

        assert isinstance(text, str)
        assert "Introducao" in text

    def test_returns_none_for_unsupported(self, doc_file: str):
        assert extract_text(doc_file) is None

    def test_returns_none_for_empty(self, pptx_empty: str):
        assert extract_text(pptx_empty) is None

    def test_returns_none_for_failed(self, corrupted_pptx: str):
        assert extract_text(corrupted_pptx) is None

    def test_legacy_caller_receives_optional_str_type_only(self, pptx_with_text: str, doc_file: str):
        """A legacy caller doing `if extract_text(...): ...` must keep working
        without ever seeing the structured ExtractionResult type."""
        for path in (pptx_with_text, doc_file):
            result = extract_text(path)
            assert result is None or isinstance(result, str)
